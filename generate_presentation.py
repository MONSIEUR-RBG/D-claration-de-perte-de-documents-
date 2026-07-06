from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

presentation = Presentation()

slides_data = [
    {
        "title": "Déclare.tg",
        "subtitle": "Plateforme nationale de déclaration de perte de documents",
        "bullets": [
            "Service officiel de la République Togolaise",
            "Déclaration de perte en ligne pour documents administratifs",
            "Suivi de dossier, validation et notifications"],
    },
    {
        "title": "Objectif du site",
        "bullets": [
            "Permettre aux citoyens de déclarer rapidement la perte de documents",
            "Simplifier le suivi administratif sans déplacement physique",
            "Centraliser les déclarations et les documents trouvés"],
    },
    {
        "title": "Utilisateurs et rôles",
        "bullets": [
            "Citoyen : déclaration de perte, suivi, consultation du dossier",
            "Agent : validation/rejet des déclarations et gestion des documents trouvés",
            "Admin : gestion des utilisateurs, types de pièces et statistiques"],
    },
    {
        "title": "Fonctionnalités principales",
        "bullets": [
            "Formulaire de déclaration de perte de document",
            "Upload de pièces justificatives (PDF, JPG, PNG)",
            "Gestion des documents trouvés et correspondances automatiques",
            "Notifications pour les utilisateurs et agents"],
    },
    {
        "title": "Flux de travail",
        "bullets": [
            "1. Déclarer une perte depuis le site",
            "2. Agent vérifie et valide/rejette la déclaration",
            "3. Document trouvé et correspondance avec une perte validée",
            "4. Notification et restitution au propriétaire"],
    },
    {
        "title": "Avantages clés",
        "bullets": [
            "Gain de temps pour le citoyen et l'administration",
            "Réduction des déplacements physiques",
            "Meilleure traçabilité des pertes et retrouvailles",
            "Plateforme sécurisée et structurée"],
    },
    {
        "title": "Technologies utilisées",
        "bullets": [
            "Laravel + PHP pour le back-end",
            "Blade templates pour les vues",
            "CSS/HTML pour l’interface utilisateur",
            "Base de données MySQL ou SQLite pour les déclarations"],
    },
]

for slide_data in slides_data:
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = slide_data["title"]
    body = content.text_frame
    body.clear()
    for idx, bullet in enumerate(slide_data["bullets"]):
        if idx == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = 'Calibri'
        p.space_after = Pt(8)

presentation.save('DeclareTogo_Presentation.pptx')
print('Presentation saved to DeclareTogo_Presentation.pptx')
